from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()

# ---------------------------
# Slide 1: Title Slide
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "DBMS: File Systems versus a DBMS"
subtitle.text = "Lecture Slide Deck with Diagram and Comparison Table"

# ---------------------------
# Slide 2: Intro Slide
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "What Is the Difference?"
body = slide.placeholders[1]
body.text = (
    "• File Systems store raw data in files\n"
    "• A DBMS manages data using structured models like tables\n"
    "• DBMS provides querying, security, concurrency, and integrity"
)

# ---------------------------
# Slide 3: File System Overview
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "File System: Overview"
slide.placeholders[1].text = (
    "• Stores data in flat files\n"
    "• No standard query language\n"
    "• Limited security and concurrency\n"
    "• Applications must handle all logic"
)

# ---------------------------
# Slide 4: DBMS Overview
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "DBMS: Overview"
slide.placeholders[1].text = (
    "• Stores data in structured tables\n"
    "• Uses SQL for powerful querying\n"
    "• Enforces ACID transactions\n"
    "• Built-in security, backups, indexing"
)

# ---------------------------
# Slide 5: Diagram
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Diagram: File System vs DBMS Architecture"

# File System Box
fs_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(4), Inches(1)
)
fs_box.text = "File System\n• Application manages structure\n• No ACID\n• No SQL"

# DBMS Box
dbms_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.5), Inches(4), Inches(1)
)
dbms_box.text = "DBMS\n• Schema + Tables\n• ACID Transactions\n• SQL Query Engine"

# Connector Arrow
line = slide.shapes.add_connector(
    1, Inches(4.5), Inches(2), Inches(5), Inches(2)
)

# ---------------------------
# Slide 6: Comparison Table
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Comparison: File System vs DBMS"

rows, cols = 8, 3
table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.2), Inches(9), Inches(4)).table

# Headings
table.cell(0, 0).text = "Feature"
table.cell(0, 1).text = "File System"
table.cell(0, 2).text = "DBMS"

data = [
    ("Data Organization", "Files", "Tables / Schema"),
    ("Querying", "Manual code", "SQL"),
    ("Redundancy", "High", "Low (Normalization)"),
    ("Concurrency", "Very limited", "Strong (Transaction control)"),
    ("Security", "File-level", "Fine-grained roles + permissions"),
    ("Backup/Recovery", "Manual", "Automatic"),
    ("Integrity", "Application-level", "Built-in constraints"),
]

for i, row in enumerate(data, start=1):
    table.cell(i, 0).text = row[0]
    table.cell(i, 1).text = row[1]
    table.cell(i, 2).text = row[2]

# ---------------------------
# Slide 7: Summary
# ---------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Summary"
slide.placeholders[1].text = (
    "• File systems provide basic data storage\n"
    "• DBMS provides structured, secure, consistent data management\n"
    "• DBMS supports indexing, transactions, backups, integrity\n"
    "• Modern applications rely heavily on DBMS systems"
)

# Save PPTX
prs.save("DBMS_vs_FileSystem.pptx")

print("PPTX file 'DBMS_vs_FileSystem.pptx' created successfully!")
