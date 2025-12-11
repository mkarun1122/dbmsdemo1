from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

# -----------------------------------------------------
# Slide 1: Title
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "DBMS: Data Independence"
slide.placeholders[1].text = "Neat diagrams and detailed content"

# -----------------------------------------------------
# Slide 2: Overview
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Overview"
slide.placeholders[1].text = (
    "• Meaning of Data Independence\n"
    "• Physical Data Independence\n"
    "• Logical Data Independence\n"
    "• Why Data Independence is important\n"
    "• Diagrams and comparison table"
)

# -----------------------------------------------------
# Slide 3: Meaning of Data Independence
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "What is Data Independence?"
slide.placeholders[1].text = (
    "• Ability to modify schema at one level without affecting higher levels\n"
    "• Achieved through a layered database architecture\n"
    "• Ensures flexibility, maintainability, and minimal application impact"
)

# -----------------------------------------------------
# Slide 4: Physical Data Independence
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Physical Data Independence"
slide.placeholders[1].text = (
    "• Internal schema changes do NOT affect conceptual schema\n"
    "• Examples of physical-level changes:\n"
    "   – File organization\n"
    "   – Index structures (B-tree, hashing)\n"
    "   – Storage formats\n"
    "   – Access paths\n"
    "• SQL queries and tables remain unchanged"
)

# -----------------------------------------------------
# Slide 5: Logical Data Independence
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Logical Data Independence"
slide.placeholders[1].text = (
    "• Ability to change conceptual schema without affecting user views\n"
    "• Examples of logical-level changes:\n"
    "   – Adding new attributes\n"
    "   – Splitting tables\n"
    "   – Adding new relationships\n"
    "• External schemas and applications remain unchanged"
)

# -----------------------------------------------------
# Slide 6: Diagram – Three-Level Architecture
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Diagram: Three-Level DBMS Architecture"

# External Level Box
ext = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(7), Inches(1)
)
ext.text = "EXTERNAL LEVEL (User Views)"

# Conceptual Level Box
conc = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(7), Inches(1)
)
conc.text = "CONCEPTUAL LEVEL (Logical Schema)"

# Internal Level Box
intr = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(7), Inches(1)
)
intr.text = "INTERNAL LEVEL (Physical Storage)"

# Arrows
slide.shapes.add_connector(1, Inches(4), Inches(2), Inches(4), Inches(2.5))
slide.shapes.add_connector(1, Inches(4), Inches(3.5), Inches(4), Inches(4))

# -----------------------------------------------------
# Slide 7: Diagram – Data Independence Flow
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Diagram: Data Independence Levels"

# Physical Change Box
phy = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(3.5), Inches(1)
)
phy.text = "Internal Schema Changes\n(Files, Indexes, Pages)"

# Logical Box
log = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(1.5), Inches(3.5), Inches(1)
)
log.text = "Conceptual Schema\n(Tables, Keys, Relations)"

# External Box
ext2 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(3.2), Inches(3.5), Inches(1)
)
ext2.text = "External Schema\n(User Views)"

# Connectors
slide.shapes.add_connector(1, Inches(2), Inches(2.5), Inches(2.5), Inches(3.2))
slide.shapes.add_connector(1, Inches(6.5), Inches(2.5), Inches(6), Inches(3.2))

# -----------------------------------------------------
# Slide 8: Table – Comparison of Data Independence
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Comparison Table: Physical vs Logical Data Independence"

rows, cols = 4, 3
table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.3), Inches(9), Inches(3)).table

# Headers
table.cell(0, 0).text = "Type"
table.cell(0, 1).text = "Lower Level Affected"
table.cell(0, 2).text = "Higher Level Unchanged"

# Rows
data = [
    ("Physical Data Independence", "Internal Schema", "Conceptual Schema"),
    ("Logical Data Independence", "Conceptual Schema", "External Schema"),
]

for i, row in enumerate(data, start=1):
    table.cell(i, 0).text = row[0]
    table.cell(i, 1).text = row[1]
    table.cell(i, 2).text = row[2]

# -----------------------------------------------------
# Slide 9: Why Data Independence Matters
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Why is Data Independence Important?"
slide.placeholders[1].text = (
    "• Reduces application maintenance cost\n"
    "• Increases database flexibility\n"
    "• Improves scalability and performance\n"
    "• Enhances security through abstraction\n"
    "• Supports long-term system evolution"
)

# -----------------------------------------------------
# Slide 10: Summary
# -----------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Summary"
slide.placeholders[1].text = (
    "• Data independence separates DB levels for flexibility\n"
    "• Physical Data Independence: internal → conceptual unaffected\n"
    "• Logical Data Independence: conceptual → external unaffected\n"
    "• Achieved through layered architecture\n"
    "• Essential for efficient and maintainable DBMS design"
)

# Save file
prs.save("DBMS_Data_Independence.pptx")
print("PowerPoint file 'DBMS_Data_Independence.pptx' created successfully!")
