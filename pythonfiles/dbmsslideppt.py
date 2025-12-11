from pptx import Presentation

prs = Presentation()

slides = [
    ("Database System Applications:\nA Historical Perspective",
     "An overview of the evolution of database systems from the 1950s to today."),

    ("Overview",
     "• Early File-Based Systems\n"
     "• Hierarchical & Network Databases\n"
     "• Relational Era\n"
     "• Web & Big Data (NoSQL)\n"
     "• Modern Data Ecosystem\n"
     "• Future Trends"),

    ("Early File-Based Systems (1950s–1960s)",
     "• Flat files (sequential, indexed)\n"
     "• High redundancy\n"
     "• Tight coupling of data & applications\n"
     "• No concurrency control"),

    ("Applications (File-Based Era)",
     "• Payroll systems\n"
     "• Inventory tracking\n"
     "• Accounting\n"
     "• Basic business operations"),

    ("Hierarchical & Network Databases (1960s–1970s)",
     "• IBM IMS (hierarchical)\n"
     "• CODASYL network model\n"
     "• Pointer-based navigation\n"
     "• Suited for stable, predictable data structures"),

    ("Relational Database Era (1970s–1990s)",
     "• Codd’s relational model\n"
     "• SQL standardization\n"
     "• Data independence\n"
     "• ACID properties\n"
     "• Widespread enterprise adoption"),

    ("Major RDBMS Platforms",
     "• Oracle\n• IBM DB2\n• Microsoft SQL Server\n• PostgreSQL\n• MySQL"),

    ("NoSQL & Big Data Era (2000s–2010s)",
     "• Need for scalability & unstructured data\n"
     "• Key-value, document, column-family, graph DBs\n"
     "• BASE vs ACID tradeoffs"),

    ("Applications (NoSQL Era)",
     "• Social media platforms\n"
     "• Search engines\n"
     "• IoT data ingestion\n"
     "• Real-time personalization\n"
     "• Massive-scale analytics"),

    ("Modern Data Ecosystem (2010s–Present)",
     "• Cloud-managed databases\n"
     "• NewSQL systems\n"
     "• Data lakes & lakehouses\n"
     "• Vector databases for AI\n"
     "• Streaming platforms (Kafka, Flink)"),

    ("Future Directions",
     "• Autonomous/self-tuning databases\n"
     "• AI-integrated semantic querying\n"
     "• Privacy-preserving computation\n"
     "• Global distributed transactional systems"),

    ("Summary",
     "Database systems evolved from simple files to relational engines, to NoSQL\n"
     "distributed systems, and now to intelligent cloud-based and AI-integrated\n"
     "platforms.")
]

for title, body in slides:
    layout = prs.slide_layouts[1]  # Title + body
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = body

prs.save("Database_System_Applications_Historical_Perspective.pptx")

print("PPTX file created successfully!")
