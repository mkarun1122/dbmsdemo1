from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# -----------------------------
# Slide 1: Title
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "DBMS: The Data Model"
slide.placeholders[1].text = "Lecture Slide Deck with Diagrams"

# -----------------------------
# Slide 2: What is a Data Model
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "What Is a Data Model?"
slide.placeholders[1].text = (
    "• A framework to describe data, relationships, and constraints\n"
    "• Defines data structure, operations, and rules\n"
    "• Essential for database design and implementation"
)

# -----------------------------
# Slide 3: Types of Data Models
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Types of Data Models"
slide.placeholders[1].text = (
    "• Relational Model\n"
    "• Entity–Relationship (ER) Model\n"
    "• Object-Based Model\n"
    "• Semi-Structured / NoSQL Models\n"
    "• Network & Hierarchical Models"
)

# ---------------------------------------------------------
# Slide 4: Diagram – Three-Level Data Modeling Architecture
# ---------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Diagram: Three-Schema Architecture"

# Boxes
concept = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(3), Inches(1)
)
concept.text = "Conceptual Level\n• High-level design\n• ER diagrams"

logical = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(1.5), Inches(3), Inches(1)
)
logical.text = "Logical Level\n• Tables, attributes, keys\n• DBMS-specific model"

physical = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(3), Inches(3), Inches(1)
)
physical.text = "Physical Level\n• Storage\n• Indexes\n• File structures"

# Arrows
slide.shapes.add_connector(1, Inches(3.8), Inches(2), Inches(4.2), Inches(2))
slide.shapes.add_connector(1, Inches(3.8), Inches(2.5), Inches(3.3), Inches(3))

# -----------------------------
# Slide 5: Relational Model Overview
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Relational Model"
slide.placeholders[1].text = (
    "• Data stored in tables (relations)\n"
    "• Rows = tuples\n"
    "• Columns = attributes\n"
    "• Keys define relationships\n"
    "• SQL used for queries"
)

# ---------------------------------------------------------
# Slide 6: Diagram – Relational Table with Keys
# ---------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Diagram: Relational Model Table"

# Table
shape = slide.shapes.add_table(4, 3, Inches(0.5), Inches(1.5), Inches(8), Inches(2)).table

# headers
shape.cell(0, 0).text = "StudentID (PK)"
shape.cell(0, 1).text = "Name"
shape.cell(0, 2).text = "DeptID (FK)"

# sample rows
shape.cell(1, 0).text = "101"
shape.cell(1, 1).text = "Alice"
shape.cell(1, 2).text = "10"

shape.cell(2, 0).text = "102"
shape.cell(2, 1).text = "John"
shape.cell(2, 2).text = "20"

shape.cell(3, 0).text = "103"
shape.cell(3, 1).text = "Maria"
shape.cell(3, 2).text = "10"

# -----------------------------
# Slide 7: ER Model Overview
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Entity–Relationship (ER) Model"
slide.placeholders[1].text = (
    "• Entity: object with attributes\n"
    "• Relationship: associations between entities\n"
    "• Graphical design tool for conceptual modeling\n"
    "• Used before creating relational schema"
)

# ---------------------------------------------------------
# Slide 8: ER Diagram (Simple)
# ---------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Diagram: ER Model"

# Entity: Student
student = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(3), Inches(1)
)
student.text = "Entity: STUDENT\n• StudentID\n• Name\n• Age"

# Entity: Department
dept = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.5), Inches(3), Inches(1)
)
dept.text = "Entity: DEPARTMENT\n• DeptID\n• DeptName"

# Relationship diamond
rel = slide.shapes.add_shape(
    MSO_SHAPE.DIAMOND, Inches(3.5), Inches(1.8), Inches(1.2), Inches(1.2)
)
rel.text = "BelongsTo"

# Connectors
slide.shapes.add_connector(1, Inches(3.2), Inches(2), Inches(3.5), Inches(2))
slide.shapes.add_connector(1, Inches(4.7), Inches(2), Inches(5), Inches(2))

# -----------------------------
# Slide 9: NoSQL Models
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Semi-Structured / NoSQL Models"
slide.placeholders[1].text = (
    "• Document model (JSON)\n"
    "• Key-value stores\n"
    "• Column-family stores\n"
    "• Graph databases\n"
    "• Schema-flexible and scalable"
)

# -----------------------------
# Slide 10: Summary
# -----------------------------
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Summary"
slide.placeholders[1].text = (
    "• Data models define how data is structured and manipulated\n"
    "• Types include relational, ER, object-oriented, and NoSQL models\n"
    "• Models include data structures, operations, and constraints\n"
    "• Foundation for database design and DBMS functionality"
)

# Save file
prs.save("DBMS_Data_Model.pptx")
print("PPTX file 'DBMS_Data_Model.pptx' created successfully!")
